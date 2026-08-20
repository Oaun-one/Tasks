import base64
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Any

from pydantic import Field, field_validator

from ..models import StrictModel
from ..source_types import (
    MediaAttachment,
    SourceCapabilityError,
    SourceCommand,
    SourceContext,
    SourceDataError,
)

RENDER_TIMEOUT_SECONDS = 60
MAX_RENDERED_SLIDES = 60
MAX_RENDERED_BYTES = 24 * 1024 * 1024


class PptxExtractTextInput(StrictModel):
    """Input for pptx.extract_text.

    Attributes:
        path: Workspace-relative PPTX path.
    """

    path: str

    @field_validator("path")
    @classmethod
    def require_pptx_extension(cls, value: str) -> str:
        """Requires the path to target a PPTX file.

        Args:
            value: Workspace-relative path from verifier.json.

        Returns:
            The unchanged path when it has a PPTX extension.

        Raises:
            ValueError: If the path does not end in ``.pptx``.
        """
        if Path(value).suffix.lower() != ".pptx":
            raise ValueError("pptx.extract_text requires a .pptx file")
        return value


class PptxExtractTextOutput(StrictModel):
    """Output from pptx.extract_text.

    Attributes:
        text: Extracted slide text.
    """

    text: str


class ExtractText(SourceCommand[PptxExtractTextInput, PptxExtractTextOutput]):
    """Extracts text from PPTX slides."""

    name = "extract_text"
    input_model = PptxExtractTextInput
    output_model = PptxExtractTextOutput

    def run(self, source_input: PptxExtractTextInput, context: SourceContext) -> PptxExtractTextOutput:
        """Runs PPTX text extraction.

        Args:
            source_input: Validated PPTX extraction input.
            context: Source runtime context.

        Returns:
            Extracted PPTX text.
        """
        resolved = context.resolve_path(source_input.path)
        return PptxExtractTextOutput(text=extract_pptx_text(resolved, context.max_content_chars))


class PptxRenderSlidesInput(StrictModel):
    """Input for ``pptx.render_slides``."""

    path: str

    @field_validator("path")
    @classmethod
    def require_pptx_extension(cls, value: str) -> str:
        if Path(value).suffix.lower() != ".pptx":
            raise ValueError("pptx.render_slides requires a .pptx file")
        return value


class RenderedSlide(StrictModel):
    """Visual evidence for one slide, with pixels excluded from JSON output."""

    number: int
    width: int
    height: int
    image_base64: str = Field(exclude=True, repr=False)


class PptxRenderSlidesOutput(StrictModel):
    """Rendered slide evidence plus ordinary extracted text."""

    text: str
    slide_count: int
    slides: list[RenderedSlide]
    truncated: bool
    renderer: str

    def media_attachments(self) -> list[MediaAttachment]:
        """Return image payloads without exposing them to reward serialization."""
        return [
            MediaAttachment(
                label=f"Slide {slide.number}",
                mime_type="image/png",
                data_base64=slide.image_base64,
            )
            for slide in self.slides
        ]


class RenderSlides(SourceCommand[PptxRenderSlidesInput, PptxRenderSlidesOutput]):
    """Render slide images for a prompt-grounded visual rubric."""

    name = "render_slides"
    input_model = PptxRenderSlidesInput
    output_model = PptxRenderSlidesOutput

    def run(
        self, source_input: PptxRenderSlidesInput, context: SourceContext
    ) -> PptxRenderSlidesOutput:
        resolved = context.resolve_path(source_input.path)
        return render_pptx_slides(resolved, context.max_content_chars)


def extract_pptx_text(path: Path, limit: int) -> str:
    """Extracts visible slide and table text from a PPTX file.

    Args:
        path: Resolved PPTX file path.
        limit: Maximum number of characters to return.

    Returns:
        Extracted text capped to ``limit`` characters.
    """
    from pptx import Presentation

    presentation = Presentation(str(path))
    parts: list[str] = []
    total = 0
    for index, slide in enumerate(presentation.slides, start=1):
        slide_parts = list(iter_slide_text(slide))
        if not slide_parts:
            continue
        total = append_limited(parts, f"[Slide {index}]", total, limit)
        if total >= limit:
            break
        for text in slide_parts:
            total = append_limited(parts, text, total, limit)
            if total >= limit:
                break
        if total >= limit:
            break
    return "\n".join(parts)[:limit]


def render_pptx_slides(path: Path, text_limit: int) -> PptxRenderSlidesOutput:
    """Render a PPTX through isolated LibreOffice and Poppler subprocesses.

    Rendering happens only in a temporary directory. The submitted deck is
    never mutated, macros are never run, and the LibreOffice user profile is
    isolated so parallel benchmark runs cannot share state.
    """
    # Parse before invoking LibreOffice so a malformed ZIP becomes the same
    # normal agent-output failure as pptx.extract_text.
    text = extract_pptx_text(path, text_limit)
    office = shutil.which("soffice") or shutil.which("libreoffice")
    pdftoppm = shutil.which("pdftoppm")
    if not office or not pdftoppm:
        raise SourceCapabilityError(
            "pptx.render_slides requires LibreOffice Impress and pdftoppm in the benchmark image"
        )

    with tempfile.TemporaryDirectory(prefix="filecheck-pptx-") as raw_temp:
        tempdir = Path(raw_temp)
        profile = tempdir / "profile"
        pdf_dir = tempdir / "pdf"
        png_dir = tempdir / "png"
        pdf_dir.mkdir()
        png_dir.mkdir()
        try:
            subprocess.run(
                [
                    office,
                    "--headless",
                    "--nologo",
                    "--nodefault",
                    "--nofirststartwizard",
                    f"-env:UserInstallation={profile.as_uri()}",
                    "--convert-to",
                    "pdf:impress_pdf_Export",
                    "--outdir",
                    str(pdf_dir),
                    str(path),
                ],
                check=True,
                capture_output=True,
                text=True,
                timeout=RENDER_TIMEOUT_SECONDS,
            )
            pdf_files = sorted(pdf_dir.glob("*.pdf"))
            if len(pdf_files) != 1:
                raise SourceDataError("PowerPoint conversion did not produce exactly one PDF")
            prefix = png_dir / "slide"
            subprocess.run(
                [pdftoppm, "-png", "-r", "110", str(pdf_files[0]), str(prefix)],
                check=True,
                capture_output=True,
                text=True,
                timeout=RENDER_TIMEOUT_SECONDS,
            )
        except subprocess.TimeoutExpired as exc:
            raise SourceCapabilityError("PPTX rendering exceeded the verifier timeout") from exc
        except subprocess.CalledProcessError as exc:
            detail = (exc.stderr or exc.stdout or "conversion failed").strip()
            raise SourceDataError(f"PowerPoint rendering failed: {detail[:300]}") from exc

        images = sorted(png_dir.glob("slide-*.png"), key=slide_sort_key)
        if not images:
            raise SourceDataError("PowerPoint rendering produced no slide images")
        slide_count = len(images)
        selected: list[RenderedSlide] = []
        total_bytes = 0
        truncated = slide_count > MAX_RENDERED_SLIDES
        for image in images[:MAX_RENDERED_SLIDES]:
            raw = image.read_bytes()
            if total_bytes + len(raw) > MAX_RENDERED_BYTES:
                truncated = True
                break
            width, height = png_dimensions(raw)
            selected.append(RenderedSlide(
                number=slide_number_from_path(image),
                width=width,
                height=height,
                image_base64=base64.b64encode(raw).decode("ascii"),
            ))
            total_bytes += len(raw)
        if not selected:
            raise SourceDataError("PowerPoint rendering exceeded the configured image evidence limit")
        return PptxRenderSlidesOutput(
            text=text,
            slide_count=slide_count,
            slides=selected,
            truncated=truncated,
            renderer="libreoffice-impress+pdftoppm",
        )


def slide_sort_key(path: Path) -> tuple[int, str]:
    """Sort Poppler's ``slide-12.png`` names numerically."""
    return slide_number_from_path(path), path.name


def slide_number_from_path(path: Path) -> int:
    """Read Poppler's trailing page number, with a stable fallback."""
    stem = path.stem.rsplit("-", 1)[-1]
    try:
        return int(stem)
    except ValueError:
        return 0


def png_dimensions(raw: bytes) -> tuple[int, int]:
    """Read PNG IHDR dimensions without adding an image-processing dependency."""
    if len(raw) < 24 or raw[:8] != b"\x89PNG\r\n\x1a\n" or raw[12:16] != b"IHDR":
        raise SourceDataError("PowerPoint renderer returned an invalid PNG")
    width = int.from_bytes(raw[16:20], "big")
    height = int.from_bytes(raw[20:24], "big")
    if width <= 0 or height <= 0:
        raise SourceDataError("PowerPoint renderer returned a PNG with invalid dimensions")
    return width, height


def iter_slide_text(slide: Any) -> list[str]:
    """Collects text from a slide's text frames and tables.

    Args:
        slide: python-pptx slide object.

    Returns:
        Non-empty text fragments from the slide.
    """
    fragments: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = shape.text.strip()
            if text:
                fragments.append(text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
                if cells:
                    fragments.append(" | ".join(cells))
    return fragments


def append_limited(parts: list[str], text: str, total: int, limit: int) -> int:
    """Appends text and returns the updated character count.

    Args:
        parts: Text fragments collected so far.
        text: Text fragment to append.
        total: Current approximate character count.
        limit: Maximum desired character count.

    Returns:
        Updated approximate character count.
    """
    parts.append(text)
    return total + len(text) + 1


COMMANDS = (ExtractText(), RenderSlides())
